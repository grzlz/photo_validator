from pptx import Presentation
import re
import glob
import random 
import datetime
from util import templates, position_dict


class Templates():
    def __init__(self, template, layout, data) -> None:
        self.template = Presentation(template)
        self.slide_layout = self.template.slide_layouts[templates.get(f'{layout}')[1]]
        self.positions = position_dict[f'{layout}']
        self.slide = self.template.slides.add_slide(self.slide_layout)
        self.data = data
        self.layout = layout

    def createTemplate(self):

        self.text_assign()

        self.pic_assign()

        self.template.save(f'{datetime.datetime.now().strftime("%Y%m%d_%H%M%S")}_{self.layout}_output.pptx')


    def text_assign(self):
        text_pos = [i for i in self.positions.keys() if re.match('text', i )]

        for place in text_pos:
            if len(self.positions[place])>1:
                for i in self.positions[place] :
                    self.slide.placeholders[i].text = self.data[place]
            else:
                self.slide.placeholders[self.positions[place][0]].text = self.data[place]

    def pic_assign(self):
        r = re.compile('pic_[0-9]+|pic_wide')
        pic_pos = [i for i in self.positions.keys() if re.match(r, i )]
        n = len(pic_pos)

        topic = [random.randint(0, len(tags) - 1) for i in range(0,n)]
        
        pics_selected = []
        for i in topic : 
            pics = [ pic for pic in pic_list if tags[i] in  pic]
            pic_selected = random.randint(0,len(pics) - 1)

            pics_selected.append(pics[pic_selected])

        for i,place in enumerate(pic_pos):
            self.slide.placeholders[self.positions[place][0]].insert_picture(pics_selected[i])


# EJEMPLO DE EJECUCION
# ppt = Templates(file, template, tex_data)
# ppt.createTemplate()

